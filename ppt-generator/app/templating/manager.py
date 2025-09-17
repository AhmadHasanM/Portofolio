from pptx import Presentation
from pptx.util import Pt, Inches, Cm
from pptx.dml.color import RGBColor
from app.branding.styler import apply_branding, style_title

def apply_template(slides_content, filepath, template: str = "general"):
    """
    Generate PPTX dengan style berbeda sesuai template.
    """
    prs = Presentation()

    if template == "modern":
        _modern_template(prs, slides_content)
    elif template == "professional":
        _professional_template(prs, slides_content)
    else:  # default = general
        _general_template(prs, slides_content)

    prs.save(filepath)


# -------------------------
# General Template
# -------------------------
def _general_template(prs, slides_content, branding="default"):
    for slide_data in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = slide_data["title"]

        tf = slide.placeholders[1].text_frame
        for point in slide_data["content"]:
            p = tf.add_paragraph()
            p.text = point
            p.font.size = Pt(18)

        # Apply branding
        style_title(slide, branding)
        apply_branding(slide, branding)


# -------------------------
# Modern Template
# -------------------------
def _modern_template(prs, slides_content):
    for i, slide_data in enumerate(slides_content):
        slide_layout = prs.slide_layouts[6]  # Blank
        slide = prs.slides.add_slide(slide_layout)

        # Title box
        title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
        tf_title = title_box.text_frame
        tf_title.text = slide_data["title"]
        tf_title.paragraphs[0].font.size = Pt(32)
        tf_title.paragraphs[0].font.bold = True
        tf_title.paragraphs[0].font.color.rgb = RGBColor(0, 102, 204)  # biru modern

        # Content box
        content_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(4))
        tf_content = content_box.text_frame
        for point in slide_data["content"]:
            p = tf_content.add_paragraph()
            p.text = f"• {point}"
            p.font.size = Pt(20)
            p.font.color.rgb = RGBColor(80, 80, 80)


# -------------------------
# Professional Template
# -------------------------
def _professional_template(prs, slides_content):
    for slide_data in slides_content:
        slide_layout = prs.slide_layouts[1]  # Title + Content
        slide = prs.slides.add_slide(slide_layout)

        # Title
        slide.shapes.title.text = slide_data["title"]
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
        slide.shapes.title.text_frame.paragraphs[0].font.bold = True
        slide.shapes.title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)  # navy

        # Content
        tf = slide.placeholders[1].text_frame
        for point in slide_data["content"]:
            p = tf.add_paragraph()
            p.text = f"- {point}"
            p.font.size = Pt(18)
            p.font.color.rgb = RGBColor(30, 30, 30)
