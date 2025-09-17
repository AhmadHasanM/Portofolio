from pptx.dml.color import RGBColor
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
import os


# 🎨 Definisi tema branding (bisa disesuaikan untuk tiap perusahaan)
BRANDING_THEMES = {
    "telkom": {
        "primary_color": RGBColor(200, 16, 46),  # merah Telkom
        "secondary_color": RGBColor(30, 30, 30),
        "font": "Arial",
        "logo_path": os.path.join("app", "branding", "assets", "telkom_logo.png")
    },
    "default": {
        "primary_color": RGBColor(0, 102, 204),
        "secondary_color": RGBColor(60, 60, 60),
        "font": "Calibri",
        "logo_path": None
    }
}


def apply_branding(slide, branding="default"):
    """
    Tambahkan branding ke slide: warna, font, logo
    """
    theme = BRANDING_THEMES.get(branding, BRANDING_THEMES["default"])

    # --- Apply background color (optional)
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(255, 255, 255)  # putih clean

    # --- Apply logo (kalau ada)
    if theme["logo_path"] and os.path.exists(theme["logo_path"]):
        slide.shapes.add_picture(theme["logo_path"], Inches(8), Inches(0.1), Inches(1.5), Inches(0.7))

    # --- Apply font styling ke semua textbox
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        for paragraph in shape.text_frame.paragraphs:
            paragraph.font.name = theme["font"]
            paragraph.font.color.rgb = theme["secondary_color"]
            if paragraph.alignment is None:
                paragraph.alignment = PP_ALIGN.LEFT


def style_title(slide, branding="default"):
    """
    Styling khusus untuk judul slide
    """
    theme = BRANDING_THEMES.get(branding, BRANDING_THEMES["default"])
    if slide.shapes.title:
        title_shape = slide.shapes.title
        if title_shape.has_text_frame:
            p = title_shape.text_frame.paragraphs[0]
            p.font.size = Pt(32)
            p.font.bold = True
            p.font.color.rgb = theme["primary_color"]
