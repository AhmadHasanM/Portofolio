import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# Folder template, logo, ilustrasi
TEMPLATE_DIR = os.path.join(os.getcwd(), "app", "templating", "templates")
LOGO_DIR = os.path.join(os.getcwd(), "app", "static", "logos")
IMAGE_DIR = os.path.join(os.getcwd(), "app", "static", "images")

# -------------------------
# Mapping template ke layout index
# -------------------------
TEMPLATE_LAYOUTS = {
    "classic": [1, 2, 3, 4, 5, 6, 7],
    "general": [0, 1, 2, 3, 4, 5, 6],
    "professional": [1, 2, 3, 4, 5, 6, 7],  # Background & Background1
}


# -------------------------
# Load template PPT fisik
# -------------------------
def load_template(template_name="general"):
    template_file = f"{template_name}_template.pptx"
    template_path = os.path.join(TEMPLATE_DIR, template_file)

    if not os.path.exists(template_path):
        template_path = os.path.join(TEMPLATE_DIR, "general_template.pptx")

    return Presentation(template_path)


# -------------------------
# Utility: Clean text
# -------------------------
def _clean_text(text: str) -> str:
    """Bersihkan format markdown seperti **bold**, *italic*, dll."""
    if not isinstance(text, str):
        return str(text)
    text = re.sub(r"\*\*(.*?)\*\*", r"\1", text)  # hapus **bold**
    text = re.sub(r"\*(.*?)\*", r"\1", text)      # hapus *italic*
    text = text.strip()
    return text


# -------------------------
# Branding helper
# -------------------------
def _apply_branding_background(slide, branding):
    background = slide.background
    fill = background.fill
    fill.solid()
    if branding == "telkom":
        fill.fore_color.rgb = RGBColor(245, 245, 245)  # light gray
    else:
        fill.fore_color.rgb = RGBColor(255, 255, 255)  # white


def _get_branding_font_color(branding):
    if branding == "telkom":
        return RGBColor(0, 51, 102)  # dark blue
    return RGBColor(30, 30, 30)  # default black


# -------------------------
# Add slide content
# -------------------------
def _add_slide_content(prs, slide_data, branding, layout_indices, template="general"):
    layout_index = layout_indices[0] if layout_indices else 0
    layout_indices.append(layout_indices.pop(0))  # rotasi layout
    slide_layout = prs.slide_layouts[layout_index]
    slide = prs.slides.add_slide(slide_layout)

    # -----------------------------
    # Business template (hapus semua placeholder bawaan)
    # -----------------------------
    if slide.shapes.title:
        slide.shapes.title.text = _clean_text(slide_data["title"])
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(30)
            paragraph.font.bold = True
            paragraph.alignment = PP_ALIGN.CENTER

        content_frame = None
        for ph in slide.placeholders:
            if ph != slide.shapes.title and hasattr(ph, "text_frame"):
                content_frame = ph.text_frame
                break

        if not content_frame:
            content_frame = slide.shapes.add_textbox(
                Inches(0.5), Inches(1.8), Inches(9), Inches(5)
            ).text_frame

        content_frame.clear()

    # -----------------------------
    # Isi konten
    # -----------------------------
    points = [_clean_text(p) for p in slide_data.get("content", [])]
    for point in points[:4]:
        para = content_frame.add_paragraph()
        para.text = str(point)
        para.font.size = Pt(22)
        para.alignment = PP_ALIGN.LEFT
        para.space_before = Pt(6)
        para.space_after = Pt(6)
        para.font.color.rgb = _get_branding_font_color(branding)

    # Logo
    logo_path = os.path.join(LOGO_DIR, f"{branding}.png")
    if os.path.exists(logo_path):
        slide.shapes.add_picture(logo_path, Inches(8), Inches(0.2), width=Inches(1.5))

    # Ilustrasi
    img_name = slide_data.get("title", "").lower().replace(" ", "_")
    for ext in [".png", ".jpg", ".jpeg"]:
        img_path = os.path.join(IMAGE_DIR, f"{img_name}{ext}")
        if os.path.exists(img_path):
            slide.shapes.add_picture(img_path, Inches(1), Inches(5.2), width=Inches(6))
            break

    _apply_branding_background(slide, branding)

    return points[4:]  # sisa konten kalau lebih dari 4


# -------------------------
# Generate PPTX file
# -------------------------
def generate_pptx_file(slides_content, output_path, template="general", branding="default"):
    os.makedirs(os.path.dirname(output_path), exist_ok=True)
    prs = load_template(template)

    layout_indices = TEMPLATE_LAYOUTS.get(template, [1, 2])

    for slide_data in slides_content:
        extra_points = _add_slide_content(prs, slide_data, branding, layout_indices)

        # Slide tambahan kalau konten lebih dari batas
        while extra_points:
            extra_slide_data = {
                "title": slide_data["title"] + " (cont.)",
                "content": extra_points
            }
            extra_points = _add_slide_content(prs, extra_slide_data, branding, layout_indices)

    prs.save(output_path)
    return output_path
